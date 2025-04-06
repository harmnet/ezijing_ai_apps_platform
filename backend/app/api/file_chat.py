from flask import Blueprint, request, jsonify
import os
import uuid
import json
import logging
import mimetypes
import datetime
import io
from werkzeug.utils import secure_filename
from app.services.llm_service import VolcanoAPI, MODEL_CONFIG, get_api_client

# 导入处理不同文件格式的库
try:
    import docx  # 用于处理.docx文件
    import PyPDF2  # 用于处理.pdf文件
    from pptx import Presentation  # 用于处理.pptx文件
    import openpyxl  # 用于处理.xlsx文件
    FILE_PROCESSORS_AVAILABLE = True
except ImportError:
    FILE_PROCESSORS_AVAILABLE = False
    logger = logging.getLogger(__name__)
    logger.warning("文件处理库未安装，某些文件类型将无法正确解析。请安装: python-docx, PyPDF2, python-pptx, openpyxl")

file_chat_bp = Blueprint('file_chat', __name__)

logger = logging.getLogger(__name__)

# 上传文件临时保存目录
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), 'uploads', 'files')

# 确保上传目录存在
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# 允许的文件类型
ALLOWED_EXTENSIONS = {'pdf', 'doc', 'docx', 'txt', 'ppt', 'pptx', 'xls', 'xlsx'}

def allowed_file(filename):
    """检查文件类型是否允许上传"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_docx(file_path):
    """从.docx文件中提取文本"""
    try:
        doc = docx.Document(file_path)
        full_text = []
        # 提取文档中的所有段落文本
        for para in doc.paragraphs:
            full_text.append(para.text)
        # 提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"从docx提取文本失败: {str(e)}")
        return f"[无法提取docx文件内容: {str(e)}]"

def extract_text_from_pdf(file_path):
    """从.pdf文件中提取文本"""
    try:
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text.append(page.extract_text())
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"从pdf提取文本失败: {str(e)}")
        return f"[无法提取pdf文件内容: {str(e)}]"

def extract_text_from_pptx(file_path):
    """从.pptx文件中提取文本"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"从pptx提取文本失败: {str(e)}")
        return f"[无法提取pptx文件内容: {str(e)}]"

def extract_text_from_xlsx(file_path):
    """从.xlsx文件中提取文本"""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"工作表: {sheet_name}")
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    value = cell.value
                    if value is not None:
                        row_values.append(str(value))
                if row_values:
                    text.append('\t'.join(row_values))
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"从xlsx提取文本失败: {str(e)}")
        return f"[无法提取xlsx文件内容: {str(e)}]"

def extract_text_from_xls(file_path):
    """从.xls文件中提取文本"""
    try:
        import xlrd  # 动态导入xlrd库处理旧格式的Excel文件
        workbook = xlrd.open_workbook(file_path)
        text = []
        
        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            sheet_name = sheet.name
            text.append(f"工作表: {sheet_name}")
            
            for row_idx in range(sheet.nrows):
                row_values = []
                for col_idx in range(sheet.ncols):
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    if cell_value is not None and cell_value != '':
                        row_values.append(str(cell_value))
                if row_values:
                    text.append('\t'.join(row_values))
        
        return '\n'.join(text)
    except ImportError:
        logger.error("xlrd库未安装，无法处理.xls文件")
        return "[无法提取xls文件内容: xlrd库未安装，请安装xlrd库]"
    except Exception as e:
        logger.error(f"从xls提取文本失败: {str(e)}")
        return f"[无法提取xls文件内容: {str(e)}]"

# 全局响应处理函数，为所有响应添加CORS头
@file_chat_bp.after_request
def add_cors_headers(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
    return response

# 处理OPTIONS请求的预检路由
@file_chat_bp.route('/chat', methods=['OPTIONS'])
@file_chat_bp.route('/upload', methods=['OPTIONS'])
def handle_options_request():
    return '', 200

@file_chat_bp.route('/upload', methods=['POST'])
def upload_file():
    """
    上传文件接口
    ---
    返回:
    {
        "success": true,
        "message": "文件上传成功",
        "file_id": "文件ID，用于后续请求"
    }
    """
    try:
        # 检查请求中是否包含文件
        if 'file' not in request.files:
            return jsonify({'success': False, 'message': '没有文件部分'}), 400
        
        file = request.files['file']
        
        # 检查文件是否为空
        if file.filename == '':
            return jsonify({'success': False, 'message': '没有选择文件'}), 400
        
        # 获取原始文件名和扩展名
        original_filename = file.filename
        file_ext = original_filename.split('.')[-1].lower() if '.' in original_filename else ''
        
        logger.info(f"接收到文件: {original_filename}, 扩展名: {file_ext}")
        
        # 检查文件类型是否允许
        if not allowed_file(original_filename):
            return jsonify({'success': False, 'message': f'文件类型不允许。允许的类型: {", ".join(ALLOWED_EXTENSIONS)}'}), 400
        
        # 安全地获取文件名并添加UUID前缀以避免文件名冲突
        filename = secure_filename(original_filename)
        file_id = str(uuid.uuid4())
        unique_filename = f"{file_id}_{filename}"
        file_path = os.path.join(UPLOAD_FOLDER, unique_filename)
        
        # 保存文件到临时目录
        file.save(file_path)
        logger.info(f"文件已保存到: {file_path}")
        
        # 获取文件大小
        file_size = os.path.getsize(file_path)
        
        # 获取文件内容类型，使用mimetypes库
        file_type = mimetypes.guess_type(file_path)[0] or 'application/octet-stream'
        
        # 提取文件文本内容，根据文件类型使用不同的处理方法
        file_content = ""
        
        # 根据文件类型选择合适的内容提取方法
        if FILE_PROCESSORS_AVAILABLE:
            try:
                if file_ext == 'txt':
                    # 对于纯文本文件，尝试不同编码
                    encodings = ['utf-8', 'gbk', 'latin-1', 'windows-1252']
                    for encoding in encodings:
                        try:
                            with open(file_path, 'r', encoding=encoding) as f:
                                file_content = f.read()
                            logger.info(f"成功使用 {encoding} 编码读取txt文件")
                            break
                        except UnicodeDecodeError:
                            continue
                    
                    if not file_content:
                        file_content = f"[无法读取txt文件内容，编码不支持]"
                
                elif file_ext == 'docx':
                    # 处理Word文档
                    file_content = extract_text_from_docx(file_path)
                    logger.info(f"成功从docx文件提取文本，长度: {len(file_content)}")
                
                elif file_ext == 'pdf':
                    # 处理PDF文档
                    file_content = extract_text_from_pdf(file_path)
                    logger.info(f"成功从pdf文件提取文本，长度: {len(file_content)}")
                
                elif file_ext == 'pptx':
                    # 处理PowerPoint文档
                    file_content = extract_text_from_pptx(file_path)
                    logger.info(f"成功从pptx文件提取文本，长度: {len(file_content)}")
                
                elif file_ext == 'xlsx':
                    # 处理Excel文档
                    file_content = extract_text_from_xlsx(file_path)
                    logger.info(f"成功从xlsx文件提取文本，长度: {len(file_content)}")
                
                elif file_ext == 'xls':
                    # 处理旧版Excel文档
                    file_content = extract_text_from_xls(file_path)
                    logger.info(f"成功从xls文件提取文本，长度: {len(file_content)}")
                
                else:
                    # 其他支持但未实现专门处理的文件类型
                    file_content = f"[不支持的文件类型: {file_ext}，请上传txt、docx、pdf、pptx、xlsx或xls文件]"
                    logger.warning(f"未实现对文件类型 {file_ext} 的专门处理")
            
            except Exception as e:
                logger.error(f"提取文件内容时发生错误: {str(e)}")
                file_content = f"[文件处理错误: {str(e)}]"
        else:
            # 文件处理库未安装，尝试以文本方式读取
            try:
                # 首先尝试常见编码
                encodings = ['utf-8', 'gbk', 'latin-1', 'windows-1252']
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            file_content = f.read()
                        logger.info(f"成功使用 {encoding} 编码读取文件")
                        break
                    except UnicodeDecodeError:
                        continue
                
                # 如果所有编码都失败
                if not file_content:
                    file_content = f"[文件处理库未安装，无法解析二进制文件: {original_filename}。请安装所需库: python-docx, PyPDF2, python-pptx, openpyxl]"
                    logger.warning(f"文件处理库未安装，无法解析二进制文件: {file_path}")
            
            except Exception as e:
                logger.error(f"读取文件内容时发生错误: {str(e)}")
                file_content = f"[文件读取错误: {str(e)}]"
        
        # 创建文件信息
        file_info = {
            'local_path': file_path,
            'original_name': original_filename,
            'file_size': file_size,
            'file_type': file_type,
            'file_content': file_content,
            'upload_time': str(datetime.datetime.now()),
            'file_extension': file_ext
        }
        
        # 将文件信息写入临时JSON文件
        with open(os.path.join(UPLOAD_FOLDER, f"{file_id}.json"), 'w', encoding='utf-8') as f:
            json.dump(file_info, f, ensure_ascii=False, indent=2)
        
        logger.info(f"文件信息已保存: {file_id}，内容长度: {len(file_content)}")
        
        return jsonify({
            'success': True,
            'message': '文件上传成功',
            'file_id': file_id,
            'file_info': {
                'name': original_filename,
                'size': file_size,
                'type': file_type
            }
        })
        
    except Exception as e:
        logger.error(f"文件上传失败: {str(e)}")
        return jsonify({'success': False, 'message': f'文件上传失败: {str(e)}'}), 500

@file_chat_bp.route('/chat', methods=['POST'])
def chat_with_file():
    """
    上传的文件对话接口
    ---
    请求示例:
    {
        "file_id": "文件ID",
        "messages": [
            {"role": "user", "content": "这个文档的内容是什么?"}
        ],
        "temperature": 0.7,
        "max_tokens": 2000
    }
    """
    try:
        data = request.json
        
        if not data:
            return jsonify({'success': False, 'message': '无效的请求数据'}), 400
        
        # 验证必需参数
        file_id = data.get('file_id')
        messages = data.get('messages', [])
        
        if not file_id:
            return jsonify({'success': False, 'message': '缺少file_id参数'}), 400
        
        if not messages or not isinstance(messages, list) or len(messages) == 0:
            return jsonify({'success': False, 'message': '消息列表为空或格式不正确'}), 400
        
        # 配置参数
        temperature = data.get('temperature', 0.7)
        max_tokens = data.get('max_tokens', 2000)
        
        # 从json文件读取文件信息
        file_info_path = os.path.join(UPLOAD_FOLDER, f"{file_id}.json")
        if not os.path.exists(file_info_path):
            return jsonify({'success': False, 'message': f'找不到文件记录: {file_id}'}), 404
        
        with open(file_info_path, 'r', encoding='utf-8') as f:
            file_info = json.load(f)
        
        # 获取文件内容
        file_content = file_info.get('file_content', '')
        file_name = file_info.get('original_name', '')
        
        # 检查文件内容是否已成功提取
        if not file_content or file_content.startswith('[无法提取') or file_content.startswith('[文件处理错误'):
            logger.warning(f"文件内容提取失败: {file_id}, {file_name}")
            # 仍然继续处理，将错误信息作为文件内容发送给模型
        else:
            logger.info(f"成功获取文件内容: {file_id}, {file_name}, 长度: {len(file_content)}")
        
        # 构建系统消息，包含文件内容
        system_message = {
            "role": "system", 
            "content": f"你是一个AI助手，正在回答用户关于以下文档的问题。请基于文档内容提供准确、有用的回答。\n\n文档名称: {file_name}\n\n文档内容:\n{file_content}"
        }
        
        # 合并系统消息和用户消息
        all_messages = [system_message] + messages
        
        try:
            # 使用get_api_client获取API客户端
            client, model_info = get_api_client("deepseek-v3-vol")
            if not client:
                logger.error(f"无法创建API客户端: {model_info.get('error')}")
                return jsonify({'success': False, 'message': f'系统配置错误: {model_info.get("error")}'}), 500
                
            logger.info(f"使用模型: deepseek-v3-vol, 调用Volcano API")
            
            # 调用大模型
            response = client.chat_completion(
                messages=all_messages,
                model_id=MODEL_CONFIG["deepseek-v3-vol"]["model_id"], 
                temperature=temperature,
                max_tokens=max_tokens
            )
            
            # 检查响应
            if 'error' not in response and 'choices' in response:
                content = response['choices'][0]['message']['content']
                logger.info(f"Volcano API响应成功: {len(content)} 字符")
                return jsonify({
                    'success': True,
                    'content': content,
                    'role': 'assistant'
                })
            else:
                error_msg = response.get('error', {}).get('message', '未知错误')
                logger.error(f"Volcano API调用失败: {error_msg}")
                return jsonify({'success': False, 'message': f'调用大模型失败: {error_msg}'}), 500
                
        except Exception as e:
            logger.error(f"调用Volcano API时发生错误: {str(e)}")
            return jsonify({'success': False, 'message': f'调用大模型时发生错误: {str(e)}'}), 500
    
    except Exception as e:
        logger.error(f"处理文件对话请求时发生错误: {str(e)}")
        return jsonify({'success': False, 'message': f'处理请求时发生错误: {str(e)}'}), 500

@file_chat_bp.route('/file_info/<file_id>', methods=['GET'])
def get_file_info(file_id):
    """
    获取已上传文件的信息和内容
    ---
    返回:
    {
        "success": true,
        "file_info": {
            "original_name": "文件名",
            "file_size": 文件大小,
            "file_type": "文件类型",
            "upload_time": "上传时间",
        },
        "content_preview": "文件内容预览（最多1000字符）",
        "content_length": 文件内容总长度
    }
    """
    try:
        # 检查文件ID是否有效
        file_info_path = os.path.join(UPLOAD_FOLDER, f"{file_id}.json")
        if not os.path.exists(file_info_path):
            return jsonify({
                'success': False, 
                'message': f'找不到文件记录: {file_id}'
            }), 404
        
        # 读取文件信息
        with open(file_info_path, 'r', encoding='utf-8') as f:
            file_info = json.load(f)
        
        # 提取文件基本信息
        file_basic_info = {
            'original_name': file_info.get('original_name', '未知'),
            'file_size': file_info.get('file_size', 0),
            'file_type': file_info.get('file_type', '未知'),
            'upload_time': file_info.get('upload_time', '未知')
        }
        
        # 提取文件内容（预览）
        file_content = file_info.get('file_content', '')
        content_length = len(file_content)
        content_preview = file_content[:1000] + '...' if content_length > 1000 else file_content
        
        # 检查内容是否成功提取
        extraction_status = 'success'
        if not file_content:
            extraction_status = 'empty'
        elif file_content.startswith('[无法提取') or file_content.startswith('[文件处理错误'):
            extraction_status = 'failed'
        
        return jsonify({
            'success': True,
            'file_info': file_basic_info,
            'content_preview': content_preview,
            'content_length': content_length,
            'extraction_status': extraction_status
        })
        
    except Exception as e:
        logger.error(f"获取文件信息失败: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'获取文件信息失败: {str(e)}'
        }), 500